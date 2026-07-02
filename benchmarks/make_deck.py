#!/usr/bin/env python3
"""Generate the Nemorix funding pitch deck as a .pptx file.

Builds a 16:9 PowerPoint deck (title, problem, insight, solution, results with
charts, differentiation, why-imec, ask, impact, appendix) with speaker notes on
every slide. Embeds the charts produced by make_charts.py.

Run make_charts.py first (this script will remind you if charts are missing).

Usage:
    python benchmarks/make_deck.py
Output:
    Nemorix_Pitch.pptx  (in the repo root)
"""
from __future__ import annotations
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Palette (imec-inspired)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x1B, 0x9A, 0xAA)
TEAL_DK = RGBColor(0x0B, 0x6E, 0x7A)
AMBER = RGBColor(0xF4, 0x9D, 0x37)
RED = RGBColor(0xD7, 0x26, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xF9)
GREY = RGBColor(0x55, 0x5F, 0x6B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)

HERE = os.path.dirname(__file__)
ROOT = os.path.abspath(os.path.join(HERE, ".."))
ASSETS = os.path.join(ROOT, "assets")

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _check_charts():
    needed = ["chart_resume_latency.png", "chart_scaling.png", "chart_cost.png",
              "chart_density.png", "chart_tiers.png"]
    missing = [n for n in needed if not os.path.exists(os.path.join(ASSETS, n))]
    if missing:
        raise SystemExit(
            "Missing charts: " + ", ".join(missing) +
            "\nRun:  python benchmarks/make_charts.py   first.")


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def add_bg(slide, color):
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=6):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(para, tuple):
            para = [para]
        for (text, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = "Segoe UI"
    return tb


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_accent_bar(slide):
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, TEAL)


def content_title(slide, kicker, title):
    add_accent_bar(slide)
    add_text(slide, Inches(0.55), Inches(0.35), Inches(12.3), Inches(0.4),
             [(kicker.upper(), 13, TEAL_DK, True, False)])
    add_text(slide, Inches(0.55), Inches(0.68), Inches(12.3), Inches(0.95),
             [(title, 30, NAVY, True, False)])


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_table(slide, x, y, w, h, headers, rows, header_color=NAVY,
              header_text=WHITE, font=12, header_font=12):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gshape.table
    # header
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = htext
        r.font.size = Pt(header_font)
        r.font.bold = True
        r.font.color.rgb = header_text
        r.font.name = "Segoe UI"
    # body
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font)
            r.font.color.rgb = DARK
            # highlight Nemorix column (last) bold teal
            if c == n_cols - 1 and ri > 0:
                r.font.bold = True
                r.font.color.rgb = TEAL_DK
            r.font.name = "Segoe UI"
    return table


def add_metric_card(slide, x, y, w, h, big, label, color=TEAL):
    add_rect(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, w, Inches(0.10), color)
    add_text(slide, x, y + Inches(0.28), w, Inches(0.85),
             [(big, 34, color, True, False)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + h - Inches(0.95), w, Inches(0.85),
             [(label, 12.5, GREY, False, False)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0)


def img(slide, name, x, y, w=None, h=None):
    path = os.path.join(ASSETS, name)
    kwargs = {}
    if w is not None:
        kwargs["width"] = w
    if h is not None:
        kwargs["height"] = h
    return slide.shapes.add_picture(path, x, y, **kwargs)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_title(prs):
    s = blank_slide(prs)
    add_bg(s, NAVY)
    add_rect(s, 0, Inches(6.95), SLIDE_W, Inches(0.55), TEAL)
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.4),
             [("Nemorix", 60, WHITE, True, False)])
    add_text(s, Inches(0.95), Inches(2.85), Inches(11.5), Inches(1.0),
             [("A Virtual Memory Runtime for Persistent AI Agents", 26, TEAL, True, False)])
    add_text(s, Inches(0.95), Inches(3.7), Inches(11.5), Inches(0.6),
             [("OS-Inspired KV-Cache Tiering:  GPU VRAM \u2192 CXL \u2192 CPU RAM \u2192 SSD",
               17, LIGHT, False, False)])
    add_text(s, Inches(0.95), Inches(5.7), Inches(11.5), Inches(0.9),
             [[("Baha eddine Ouni", 16, WHITE, True, False)],
              [("imec Innovation Challenge \u00b7 2026  \u00b7  Ask: new internal R&D track",
                13, LIGHT, False, False)]])
    set_notes(s, "Persistent AI agents are coming to every enterprise. The thing stopping "
                 "them isn't compute -- it's memory. Nemorix is an operating system for agent "
                 "memory. In simulation, it runs 50 agents on a GPU that fits 6, wakes them 120 "
                 "times faster, and cuts cost per agent by about 85%. Let me show you.")


def slide_problem(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Problem", "AI agents lose their memory when the GPU fills up")
    bullets = [
        ("One long-context agent = ~20 GB of KV-cache. An 80 GB H100 fits only 4\u20136 agents.", ),
        ("When VRAM fills, today's engines discard agent context and recompute from scratch.", ),
        ("Recompute cost: 1\u20132 seconds of latency on every single wake-up.", ),
        ("Result: low agent density, slow resumes, high cost.", ),
    ]
    runs = []
    for b in bullets:
        runs.append([("\u25B6  ", 16, TEAL, True, False), (b[0], 17, DARK, False, False)])
    add_text(s, Inches(0.7), Inches(2.0), Inches(7.3), Inches(4.0), runs,
             line_spacing=1.15, space_after=18)
    # right callout
    add_rect(s, Inches(8.4), Inches(2.0), Inches(4.2), Inches(3.4), NAVY)
    add_text(s, Inches(8.6), Inches(2.35), Inches(3.8), Inches(3.0),
             [[("1\u20132 s", 54, AMBER, True, False)],
              [("wasted on every", 16, WHITE, False, False)],
              [("agent resume", 16, WHITE, False, False)],
              [("", 8, WHITE, False, False)],
              [("0 of 50 agents", 22, RED, True, False)],
              [("meet a 200 ms SLA today", 14, LIGHT, False, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s, "When you run an AI agent, it builds up a memory called the KV-cache -- for a "
                 "long conversation that's about 20 gigabytes. A top-end H100 GPU only fits a "
                 "handful. So today's engines -- vLLM, TensorRT -- do the brutal thing: when memory "
                 "fills, they throw the agent's memory away. Next time it wakes, it recomputes from "
                 "scratch. One to two seconds of dead time, every resume. It caps how many agents "
                 "you can run, and it's expensive.")


def slide_timing(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "Why Now", "The bottleneck isn't compute \u2014 it's memory management")
    add_text(s, Inches(0.7), Inches(1.95), Inches(12.0), Inches(0.8),
             [[("The bottleneck in agentic AI is the absence of a memory layer built for "
                "agent persistence.", 18, NAVY, True, False)]], line_spacing=1.1)
    cards = [
        ("75\u201398%", "of an agent's life is spent idle\n(waiting on humans, tools, CI)", AMBER),
        ("100s", "of concurrent agents enterprises\nwant \u2014 not six", TEAL),
        ("2024\u201326", "shift from stateless chat to\npersistent, long-running agents", TEAL_DK),
    ]
    x = Inches(0.7)
    for big, label, color in cards:
        add_metric_card(s, x, Inches(3.0), Inches(3.85), Inches(2.7), big, label, color)
        x += Inches(4.1)
    set_notes(s, "The market just shifted from stateless chatbots to persistent agents that run "
                 "for hours. Those agents are idle 75 to 98 percent of the time -- yet their memory "
                 "sits in the most expensive real estate in the data center. Nobody has built a "
                 "memory manager designed for this. That's the gap.")


def slide_insight(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Insight", "Agents are processes. KV-cache is virtual memory.")
    add_table(
        s, Inches(0.7), Inches(2.0), Inches(7.4), Inches(3.6),
        ["Operating System", "Nemorix"],
        [["Process", "AI agent"],
         ["RAM (fast, small)", "GPU VRAM"],
         ["Swap / disk (slow, big)", "CXL \u2192 RAM \u2192 SSD"],
         ["Page fault", "Agent wakes, cache not in GPU"],
         ["Page replacement (LRU)", "Semantic eviction"]],
        font=14, header_font=14)
    add_rect(s, Inches(8.5), Inches(2.0), Inches(4.1), Inches(3.6), TEAL)
    add_text(s, Inches(8.7), Inches(2.4), Inches(3.7), Inches(2.9),
             [[("Page it,", 30, WHITE, True, False)],
              [("don't delete it.", 30, WHITE, True, False)],
              [("", 10, WHITE, False, False)],
              [("60 years of proven OS theory, "
                "applied to AI agents for the first time.", 14, LIGHT, False, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    set_notes(s, "Operating systems solved this in the 1960s. Your laptop runs more apps than "
                 "fit in RAM by paging idle parts out to disk and back -- invisibly. We do exactly "
                 "that for agents: an agent is a process, its KV-cache is its memory, and instead of "
                 "deleting it when the GPU is full, we page it down to cheaper storage and bring it "
                 "back in milliseconds. The idea stands on 60 years of proven theory.")


def slide_solution(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Solution", "A 4-tier memory runtime for agent KV-cache")
    img(s, "chart_tiers.png", Inches(0.5), Inches(1.95), w=Inches(7.4))
    add_text(s, Inches(8.1), Inches(2.05), Inches(4.6), Inches(4.6),
             [[("Plus four mechanisms:", 16, NAVY, True, False)],
              [("\u2022 Lifecycle scheduling", 15, DARK, False, False)],
              [("   move idle agents down proactively", 12, GREY, False, False)],
              [("\u2022 Semantic eviction", 15, DARK, False, False)],
              [("   keep what matters most", 12, GREY, False, False)],
              [("\u2022 Progressive page-in", 15, DARK, False, False)],
              [("   load first layers, stream the rest", 12, GREY, False, False)],
              [("\u2022 Compression FP16\u2192FP8\u2192INT4", 15, DARK, False, False)],
              [("   per-tier, <1.8% quality loss", 12, GREY, False, False)],
              [("", 8, GREY, False, False)],
              [("CXL is the enabler: 10\u00d7 cheaper than", 13, TEAL_DK, True, False)],
              [("VRAM, fast enough for <100 ms resume.", 13, TEAL_DK, True, False)]],
             line_spacing=1.05, space_after=4)
    set_notes(s, "Nemorix manages each agent's memory across four tiers. The star is CXL -- "
                 "Compute Express Link -- 10x cheaper than GPU memory but fast enough to wake an "
                 "agent in under 100 milliseconds. We add lifecycle-aware scheduling, semantic "
                 "eviction that decides what's worth keeping, and progressive page-in that loads "
                 "just the first layers an agent needs to start responding.")


def slide_results(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Results \u2605", "120\u00d7 faster resume \u2014 verified across 8 seeds")
    img(s, "chart_resume_latency.png", Inches(0.5), Inches(1.95), w=Inches(7.3))
    # metric cards on the right
    add_metric_card(s, Inches(8.1), Inches(1.95), Inches(2.25), Inches(1.7),
                    "120\u00d7", "faster resume\n(1,151 \u2192 9.5 ms)", TEAL)
    add_metric_card(s, Inches(10.5), Inches(1.95), Inches(2.25), Inches(1.7),
                    "50/50", "agents under SLA\n(vs 0/50 today)", TEAL_DK)
    add_metric_card(s, Inches(8.1), Inches(3.85), Inches(2.25), Inches(1.7),
                    "~85%", "lower cost\nper agent-hour", AMBER)
    add_metric_card(s, Inches(10.5), Inches(3.85), Inches(2.25), Inches(1.7),
                    "8\u00d7", "agent density\n(50 vs 6 per GPU)", TEAL)
    add_text(s, Inches(8.1), Inches(5.75), Inches(4.65), Inches(0.7),
             [[("Spread across seeds \u2248 2% \u2014 not a cherry-picked result.",
                12.5, GREY, False, True)]], line_spacing=1.0)
    set_notes(s, "Here's the proof. Average resume latency drops from about 1,150 milliseconds to "
                 "9.5 -- 120x faster. Every one of the 50 agents now meets a 200-millisecond SLA; "
                 "with recompute, zero do. Cost per agent-hour falls about 85%. And these are "
                 "averaged over eight random seeds with only about 2% variation. Not one lucky run.")


def slide_scaling(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Differentiator \u2605", "Semantic eviction wins under pressure")
    img(s, "chart_scaling.png", Inches(0.5), Inches(1.9), w=Inches(7.6))
    add_table(
        s, Inches(8.3), Inches(2.6), Inches(4.4), Inches(1.6),
        ["Policy @ 500 agents", "SLA"],
        [["LRU (recency only)", "65 / 500"],
         ["Nemorix (semantic)", "273 / 500"]],
        font=14, header_font=13)
    add_text(s, Inches(8.3), Inches(4.5), Inches(4.4), Inches(2.2),
             [[("4.2\u00d7 more agents", 22, TEAL_DK, True, False)],
              [("kept under SLA at scale.", 16, DARK, False, False)],
              [("", 8, DARK, False, False)],
              [("Protects high-priority, high-attention, "
                "expensive-to-rebuild agents \u2014 the defensible "
                "algorithmic edge.", 13, GREY, False, False)]],
             line_spacing=1.1)
    set_notes(s, "At small scale, LRU and Nemorix are similar -- CXL does the heavy lifting. The "
                 "difference shows at 500 agents when the hierarchy is full. Plain LRU keeps 65 "
                 "agents under SLA. Nemorix keeps 273 -- over four times more -- because semantic "
                 "eviction protects the agents that matter. That's our defensible algorithmic edge.")


def slide_sota(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "Beyond State of the Art", "First to combine CXL + lifecycle + semantic")
    add_table(
        s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.4),
        ["System", "Multi-tier", "CXL", "Semantic eviction", "Agent-centric"],
        [["vLLM", "partial", "No", "No (LRU)", "No (request)"],
         ["NVIDIA Dynamo", "Yes", "No", "No", "semi"],
         ["FlexGen / Mooncake", "Yes", "No", "No", "No (throughput)"],
         ["Liaw & Chen '25", "\u2014", "Yes (training)", "\u2014", "No"],
         ["Nemorix", "Yes (4-tier)", "Yes", "Yes", "Yes"]],
        font=13, header_font=12)
    add_text(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.8),
             [[("Every ingredient exists somewhere. Nobody combines them. Nemorix is the only "
                "system at the intersection.", 15, NAVY, True, False)]], line_spacing=1.1)
    set_notes(s, "Every ingredient exists, but nobody combines them. vLLM is request-centric and "
                 "LRU-only. Dynamo has tiers but no CXL. FlexGen and Mooncake optimize batch "
                 "throughput, not agent persistence. There's CXL work -- but for training, a "
                 "different problem. Nemorix is the only system at the intersection of CXL, "
                 "agent-lifecycle scheduling, and semantic eviction.")


def slide_why_imec(prs):
    s = blank_slide(prs)
    add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, TEAL)
    add_text(s, Inches(0.55), Inches(0.35), Inches(12.3), Inches(0.4),
             [("WHY IMEC", 13, TEAL, True, False)])
    add_text(s, Inches(0.55), Inches(0.68), Inches(12.3), Inches(0.95),
             [("Uniquely positioned: simulation \u2192 silicon", 30, WHITE, True, False)])
    items = [
        "World-class CXL & memory-system expertise",
        "Hardware-software co-design capability",
        "Access to semiconductor partners (Samsung, SK Hynix)",
        "HPC infrastructure for real H100 + CXL validation",
        "Deep-tech commercialization support",
    ]
    runs = [[("\u25B6  ", 16, TEAL, True, False), (it, 18, WHITE, False, False)] for it in items]
    add_text(s, Inches(0.8), Inches(2.2), Inches(11.0), Inches(4.0), runs,
             line_spacing=1.2, space_after=16)
    set_notes(s, "This is why imec is the right home. Nemorix's bet is CXL -- and imec has "
                 "world-class CXL and memory expertise, hardware-software co-design, direct "
                 "relationships with Samsung and SK Hynix, and the HPC infrastructure to validate "
                 "on real hardware. Almost no one else on earth can take this from simulation to "
                 "production. That's our unfair advantage.")


def slide_business_model(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Business Model", "Open core \u2014 the playbook behind $62B Databricks")
    # Three product tiers
    add_table(
        s, Inches(0.7), Inches(2.0), Inches(7.6), Inches(2.6),
        ["Product", "Who buys it", "Price"],
        [["Nemorix Core (OSS)", "Everyone \u2014 adoption engine", "Free (MIT license)"],
         ["Nemorix Enterprise", "Cloud providers, AI companies", "$15\u201325K / GPU / year"],
         ["Nemorix Cloud", "Startups, mid-market", "$0.02 / agent-hour"]],
        font=13, header_font=12)
    # ROI callout
    add_rect(s, Inches(8.6), Inches(2.0), Inches(4.0), Inches(2.6), NAVY)
    add_text(s, Inches(8.8), Inches(2.3), Inches(3.6), Inches(2.0),
             [[("10:1 ROI", 34, AMBER, True, False)],
              [("", 6, WHITE, False, False)],
              [("Customer saves $150K/yr", 14, WHITE, False, False)],
              [("in GPU costs, pays us $15K", 14, WHITE, False, False)],
              [("", 6, WHITE, False, False)],
              [("Easy CFO approval.", 13, TEAL, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    # Revenue projection cards
    cards = [
        ("Year 2", "$750K ARR", "50 customers"),
        ("Year 3", "$5M ARR", "200 customers"),
        ("Year 5", "$100M ARR", "2,000 customers"),
    ]
    x = Inches(0.7)
    for yr, arr, cust in cards:
        add_rect(s, x, Inches(5.0), Inches(3.85), Inches(1.8), LIGHT)
        add_rect(s, x, Inches(5.0), Inches(3.85), Inches(0.10), TEAL)
        add_text(s, x, Inches(5.2), Inches(3.85), Inches(0.5),
                 [(yr, 13, GREY, True, False)], align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(5.55), Inches(3.85), Inches(0.7),
                 [(arr, 26, TEAL_DK, True, False)], align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(6.15), Inches(3.85), Inches(0.5),
                 [(cust, 12, GREY, False, False)], align=PP_ALIGN.CENTER)
        x += Inches(4.1)
    set_notes(s, "How does this make money? Same playbook that built Databricks into a 62-billion "
                 "dollar company: open core. The base runtime is free and open source -- that drives "
                 "adoption. The enterprise version adds multi-node CXL pooling, dashboards, RBAC, "
                 "and support at 15 to 25 thousand per GPU per year. And a managed cloud service for "
                 "pay-per-use. A customer saving 150K a year in GPU costs happily pays us 15K -- "
                 "that's 10-to-1 ROI. Year 5 target: 100 million ARR, 85% gross margin.")


def slide_market_exit(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "Market & Exit", "$45B market \u2014 three paths to outsized returns")
    # TAM / SAM / SOM cards
    cards = [
        ("$45B", "TAM", "AI inference\ninfrastructure", NAVY),
        ("$8B", "SAM", "Persistent agent\nserving", TEAL_DK),
        ("$100M+", "SOM (Yr 5)", "Our realistic\ncapture", TEAL),
    ]
    x = Inches(0.7)
    for big, label, desc, color in cards:
        add_rect(s, x, Inches(2.0), Inches(3.0), Inches(2.0), color)
        add_text(s, x, Inches(2.2), Inches(3.0), Inches(0.7),
                 [(big, 30, WHITE, True, False)], align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.85), Inches(3.0), Inches(0.4),
                 [(label, 13, AMBER, True, False)], align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.2), Inches(3.0), Inches(0.7),
                 [(desc, 11, LIGHT, False, False)], align=PP_ALIGN.CENTER)
        x += Inches(3.25)
    # Exit paths with comparables
    add_rect(s, Inches(10.5), Inches(2.0), Inches(2.15), Inches(2.0), AMBER)
    add_text(s, Inches(10.5), Inches(2.2), Inches(2.15), Inches(1.8),
             [[("25\u201335%", 24, NAVY, True, False)],
              [("CAGR", 13, NAVY, True, False)],
              [("", 4, NAVY, False, False)],
              [("Growing every", 11, NAVY, False, False)],
              [("quarter", 11, NAVY, False, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Exit path table
    add_table(
        s, Inches(0.7), Inches(4.4), Inches(11.9), Inches(2.4),
        ["Exit Path", "How", "Potential"],
        [["Category leader", "\"Nemorix\" = agent memory (like K8s = orchestration)", "$2B+"],
         ["Strategic acquisition", "NVIDIA, Samsung, Intel buy the software layer", "$500M\u2013$2B"],
         ["Hardware partnership", "CXL vendors co-invest to drive adoption", "Revenue + equity"]],
        font=13, header_font=12)
    add_text(s, Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.5),
             [[("Comparable exits: Databricks $62B \u00b7 HashiCorp $6.8B (IBM) \u00b7 "
                "Mellanox $6.9B (NVIDIA) \u00b7 Redis $2B+",
                12, GREY, False, True)]], line_spacing=1.0)
    set_notes(s, "The AI inference infrastructure market is 45 billion, growing 25 to 35 percent "
                 "a year. Persistent agent serving is an 8 billion segment. Three exit paths: "
                 "become the category leader like Kubernetes, strategic acquisition by NVIDIA or "
                 "Samsung, or hardware partnerships. For imec: equity, IP, and recognition as the "
                 "birthplace of AI virtual memory. Comparable exits range from 2 to 62 billion.")


def slide_ask(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "The Ask & Plan", "Fund a new R&D track to validate on real hardware")
    add_table(
        s, Inches(0.7), Inches(2.0), Inches(8.4), Inches(3.6),
        ["Phase", "What", "Outcome"],
        [["Now \u2713", "Simulator + paper draft", "74 tests, preprint-ready"],
         ["0\u20136 mo", "Workshop paper + arXiv", "priority + credibility"],
         ["3\u20139 mo", "vLLM + real H100 + CXL DIMM", "measured, not simulated"],
         ["6\u201312 mo", "Eviction study + OSS release", "adoption + citations"],
         ["12\u201336 mo", "Spin-out / license / upstream", "value capture"]],
        font=12.5, header_font=12.5)
    add_rect(s, Inches(9.4), Inches(2.0), Inches(3.2), Inches(3.6), TEAL)
    add_text(s, Inches(9.6), Inches(2.4), Inches(2.8), Inches(3.0),
             [[("\u20ac80\u2013120K", 28, WHITE, True, False)],
              [("hardware POC", 15, LIGHT, False, False)],
              [("", 8, WHITE, False, False)],
              [("CXL module\n+ engineering time\n+ cloud GPU credits", 14, WHITE, False, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    set_notes(s, "What I'm asking for is 80 to 120 thousand euros to validate Nemorix on real "
                 "hardware. The simulator is done. The next milestone is integrating with vLLM and "
                 "running on a real H100 with a CXL DIMM here at imec. From there: a paper, an "
                 "open-source release, and a choice of spin-out or license.")


def slide_impact(prs):
    s = blank_slide(prs)
    add_bg(s, NAVY)
    add_rect(s, 0, Inches(6.95), SLIDE_W, Inches(0.55), TEAL)
    add_text(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.6),
             [[("Turn the GPU memory wall", 36, WHITE, True, False)],
              [("into a memory hierarchy.", 36, TEAL, True, False)]], line_spacing=1.05)
    items = [
        "Lower resume latency \u2192 real-time, long-context agents",
        "Lower cost & higher density \u2192 labs, hospitals, public institutions can deploy",
        "More sustainable use of existing GPUs",
    ]
    runs = [[("\u25B6  ", 16, TEAL, True, False), (it, 17, LIGHT, False, False)] for it in items]
    add_text(s, Inches(0.95), Inches(3.5), Inches(11.0), Inches(2.2), runs,
             line_spacing=1.2, space_after=14)
    add_text(s, Inches(0.95), Inches(5.9), Inches(11.0), Inches(0.7),
             [[("Make persistent AI agents practical without hyperscale infrastructure.",
                16, WHITE, True, True)]])
    set_notes(s, "The impact is democratization. Today, running fleets of persistent agents needs "
                 "hyperscale budgets. Nemorix lets a research lab, a hospital, or a public "
                 "institution do it on the hardware they already have -- faster, cheaper, more "
                 "sustainably. We turn the GPU memory wall into a memory hierarchy. Thank you -- "
                 "I'd love your questions.")


def slide_appendix_density_cost(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "Appendix", "Supporting data: density & cost")
    img(s, "chart_density.png", Inches(0.6), Inches(2.0), w=Inches(5.9))
    img(s, "chart_cost.png", Inches(6.8), Inches(2.0), w=Inches(5.9))
    set_notes(s, "Backup slide. Left: Nemorix raises concurrent agents under SLA from 6 to 50 on "
                 "one GPU. Right: cost per agent-hour drops about 85%. Both reproducible from "
                 "results.json.")


def slide_appendix_repro(prs):
    s = blank_slide(prs)
    add_bg(s, WHITE)
    content_title(s, "Appendix", "Reproducibility & honesty")
    add_rect(s, Inches(0.7), Inches(2.0), Inches(7.2), Inches(3.4), NAVY)
    add_text(s, Inches(0.95), Inches(2.25), Inches(6.7), Inches(3.0),
             [[("Reproduce every number:", 15, TEAL, True, False)],
              [("pip install -e \".[dev]\"", 13, WHITE, False, False)],
              [("python -m pytest tests/ -v", 13, WHITE, False, False)],
              [("python benchmarks/run_simulation.py", 13, WHITE, False, False)],
              [("python benchmarks/run_robustness.py", 13, WHITE, False, False)],
              [("python benchmarks/compare_policies.py", 13, WHITE, False, False)],
              [("", 8, WHITE, False, False)],
              [("74 automated tests \u00b7 every constant from a datasheet",
                13, LIGHT, False, False)]],
             line_spacing=1.25)
    add_text(s, Inches(8.2), Inches(2.0), Inches(4.5), Inches(4.5),
             [[("The honesty line:", 15, NAVY, True, False)],
              [("\u201cThis is simulation, calibrated to "
                "published hardware specs \u2014 not yet measured "
                "on real hardware. That's exactly what the "
                "hardware-validation phase is for.\u201d", 15, DARK, False, True)],
              [("", 8, DARK, False, False)],
              [("Sources: H100 HBM3 datasheet, Samsung CMM-D "
                "CXL 2.0, NVMe Gen4, MLPerf v4.0, KIVI/KVQuant.",
                12, GREY, False, False)]],
             line_spacing=1.2)
    set_notes(s, "Backup. Every number is reproducible with these commands, and every hardware "
                 "constant traces to a published datasheet. Always lead with the honesty line: this "
                 "is simulation calibrated to specs, and hardware validation is the next milestone.")


def main():
    _check_charts()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_timing(prs)
    slide_insight(prs)
    slide_solution(prs)
    slide_results(prs)
    slide_scaling(prs)
    slide_sota(prs)
    slide_why_imec(prs)
    slide_business_model(prs)
    slide_market_exit(prs)
    slide_ask(prs)
    slide_impact(prs)
    slide_appendix_density_cost(prs)
    slide_appendix_repro(prs)

    out = os.path.join(ROOT, "Nemorix_Pitch.pptx")
    # If file is locked (e.g. open in PowerPoint), save to alternative name
    try:
        prs.save(out)
    except PermissionError:
        out = os.path.join(ROOT, "Nemorix_Pitch_v2.pptx")
        prs.save(out)
    print(f"Saved deck: {os.path.relpath(out)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
